"""PowerPoint presentation renderer using python-pptx."""

from typing import Any
from uuid import uuid4
from pathlib import Path
import logging
import io
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import httpx

logger = logging.getLogger(__name__)


class PowerPointRenderer:
    """Renders PowerPoint presentations from Layout Specification Packages."""

    def __init__(self, output_dir: str | None = None) -> None:
        """
        Initialize PowerPoint renderer.

        Args:
            output_dir: Directory to save generated presentations (defaults to infrastructure/data/artifacts)
        """
        if output_dir is None:
            # Use absolute path relative to project root
            from pathlib import Path
            project_root = Path(__file__).parent.parent.parent.parent.resolve()
            output_dir = project_root / "infrastructure" / "data" / "artifacts"
        self.output_dir = Path(output_dir).resolve()
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self._http_client = None  # Lazy initialization for image downloads

    async def render(self, lsp: dict[str, Any]) -> str:
        """
        Render PowerPoint presentation from LSP.

        Args:
            lsp: Layout Specification Package

        Returns:
            Artifact ID for the generated presentation
        """
        # Extract metadata
        metadata = lsp.get("metadata", {})
        structure = lsp.get("structure", [])
        proposal_id = lsp.get("proposal_id", "unknown")
        session_id = lsp.get("session_id", "unknown")

        # Create presentation
        prs = Presentation()

        # Set presentation dimensions to 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        title = metadata.get("title", "Untitled Presentation")

        logger.info(f"Rendering PowerPoint presentation: {title} with {len(structure)} slides")
        logger.debug(f"Content map has {len(lsp.get('content_map', {}))} entries: {list(lsp.get('content_map', {}).keys())[:5]}")
        logger.debug(f"First structure unit has {len(structure[0].get('elements', [])) if structure else 0} elements")

        # Iterate through structure units (slides)
        for unit_idx, unit in enumerate(structure):
            unit_type = unit.get("type", "slide")
            unit_title = unit.get("title")
            template = unit.get("template", "standard_content")
            elements = unit.get("elements", [])
            notes = unit.get("notes")

            # Select slide layout based on template
            slide_layout = self._select_layout(prs, template, unit_idx == 0)

            # Add slide
            slide = prs.slides.add_slide(slide_layout)

            # Render elements on slide
            await self._render_slide_elements(slide, elements, unit_title, lsp)

            # Add speaker notes if present
            if notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes

        # Generate artifact ID and save
        artifact_id = f"artifact-pptx-{uuid4().hex[:8]}"
        output_path = self.output_dir / f"{artifact_id}.pptx"

        prs.save(str(output_path))

        logger.info(f"PowerPoint presentation saved: {output_path} (artifact_id: {artifact_id})")

        return artifact_id

    def _select_layout(self, prs: Presentation, template: str, is_title: bool):
        """
        Select appropriate slide layout.

        Args:
            prs: Presentation instance
            template: Template name from LSP
            is_title: Whether this is the title slide

        Returns:
            Slide layout
        """
        # python-pptx default layouts:
        # 0 = Title Slide
        # 1 = Title and Content
        # 5 = Title Only
        # 6 = Blank

        if is_title or template == "title_slide":
            return prs.slide_layouts[0]  # Title Slide
        elif template in ["bullet_list", "standard_content"]:
            return prs.slide_layouts[1]  # Title and Content
        elif template == "title_only":
            return prs.slide_layouts[5]  # Title Only
        else:
            return prs.slide_layouts[6]  # Blank

    async def _render_slide_elements(
        self, slide, elements: list[dict[str, Any]], slide_title: str | None, lsp: dict[str, Any]
    ) -> None:
        """
        Render elements on a slide.

        Args:
            slide: Slide object
            elements: List of layout elements
            slide_title: Optional slide title
            lsp: Full LSP for content resolution
        """
        # Add title if present
        if slide_title and slide.shapes.title:
            slide.shapes.title.text = slide_title

        # Render each element based on its type
        for idx, element in enumerate(elements):
            element_type = element.get("type", "text")
            content_ref = element.get("content_ref")
            position = element.get("position", {})
            styling = element.get("styling", {})
            gestalt_rules = element.get("gestalt_rules", {})

            # Resolve content
            content_text = self._resolve_content(content_ref, lsp)

            if not content_text or content_text.startswith("[Missing content"):
                logger.warning(f"Skipping element with missing content: content_ref={content_ref}, "
                             f"content_map_keys={list(lsp.get('content_map', {}).keys())[:5]}")
                continue

            # Route to appropriate renderer based on element type
            if element_type == "text":
                await self._render_text_element(slide, element, content_text, position, styling, gestalt_rules, idx)
            elif element_type == "image":
                await self._render_image_element(slide, element, content_text, position, styling)
            elif element_type == "table":
                await self._render_table_element(slide, element, content_text, position, styling)
            elif element_type == "shape":
                self._render_shape_element(slide, element, position, styling)
            else:
                logger.warning(f"Unsupported element type: {element_type}")

    def _resolve_content(self, content_ref: str | None, lsp: dict[str, Any]) -> str | None:
        """Resolve content text from content_ref."""
        if not content_ref:
            return None

        # Look up content in content_map
        content_map = lsp.get("content_map", {})
        return content_map.get(content_ref, f"[Missing content: {content_ref}]")

    def _hex_to_rgb(self, hex_color: str) -> tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    async def _render_text_element(
        self,
        slide,
        element: dict[str, Any],
        content_text: str,
        position: dict[str, Any],
        styling: dict[str, Any],
        gestalt_rules: dict[str, Any],
        idx: int,
    ) -> None:
        """
        Render a text element with support for bullets and multi-line text.

        Args:
            slide: Slide object
            element: Full element configuration
            content_text: Resolved content text
            position: Position configuration
            styling: Styling configuration
            gestalt_rules: Gestalt rules
            idx: Element index for default positioning
        """
        # Get position (convert from LSP inches to python-pptx)
        left = Inches(position.get("x", 1.0))
        top = Inches(position.get("y", 1.5 + idx * 0.8))
        width = Inches(position.get("width", 8.5))

        # Determine if this is a bullet list
        is_bullet_list = self._is_bullet_content(content_text)

        # Calculate height based on content
        base_height = position.get("height", 0.5)
        if is_bullet_list or len(content_text) > 200:
            # Estimate height based on lines
            lines = content_text.count('\n') + 1 if '\n' in content_text else max(1, len(content_text) // 60)
            height = Inches(min(base_height * lines * 0.3, 5.0))  # Max 5 inches
        else:
            height = Inches(base_height)

        # Add text box with word wrap enabled
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Apply background color to the shape if specified
        background = styling.get("background")
        if background:
            fill = textbox.fill
            fill.solid()
            rgb = self._hex_to_rgb(background)
            fill.fore_color.rgb = RGBColor(*rgb)

        # Handle bullet lists
        if is_bullet_list:
            self._format_as_bullets(text_frame, content_text, styling, gestalt_rules)
        else:
            # Handle multi-paragraph text
            paragraphs = content_text.split('\n\n')
            for para_idx, para_text in enumerate(paragraphs):
                if para_idx == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()

                paragraph.text = para_text.strip()
                self._apply_paragraph_styling(paragraph, styling, gestalt_rules)

    def _is_bullet_content(self, text: str) -> bool:
        """
        Determine if content should be formatted as bullets.

        Args:
            text: Content text

        Returns:
            True if content appears to be a bullet list
        """
        # Check for common bullet point patterns
        bullet_patterns = [
            r'^\s*[-•*]\s+',  # Starts with -, •, or *
            r'^\s*\d+\.\s+',  # Numbered list
            r'\n\s*[-•*]\s+',  # Contains bullet points
            r'\n\s*\d+\.\s+',  # Contains numbered items
        ]

        for pattern in bullet_patterns:
            if re.search(pattern, text, re.MULTILINE):
                return True

        return False

    def _format_as_bullets(
        self,
        text_frame,
        content_text: str,
        styling: dict[str, Any],
        gestalt_rules: dict[str, Any],
    ) -> None:
        """
        Format text as bullet points.

        Args:
            text_frame: TextFrame object
            content_text: Content text with bullet markers
            styling: Styling configuration
            gestalt_rules: Gestalt rules
        """
        # Split into bullet points (handle various formats)
        lines = content_text.split('\n')
        bullet_items = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Remove bullet markers
            line = re.sub(r'^\s*[-•*]\s+', '', line)
            line = re.sub(r'^\s*\d+\.\s+', '', line)

            if line:
                bullet_items.append(line)

        # Add bullet points to text frame
        for idx, item in enumerate(bullet_items):
            if idx == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()

            paragraph.text = item
            paragraph.level = 0  # Top-level bullet

            # Apply styling
            self._apply_paragraph_styling(paragraph, styling, gestalt_rules)

    def _apply_paragraph_styling(
        self,
        paragraph,
        styling: dict[str, Any],
        gestalt_rules: dict[str, Any],
    ) -> None:
        """
        Apply styling to a paragraph.

        Args:
            paragraph: Paragraph object
            styling: Styling configuration
            gestalt_rules: Gestalt rules
        """
        # Font styling
        font_config = styling.get("font", {})
        if font_config and paragraph.runs:
            font = paragraph.runs[0].font
            if "size_pt" in font_config:
                font.size = Pt(font_config["size_pt"])
            if "family" in font_config:
                font.name = font_config["family"]
            if "weight" in font_config and font_config["weight"] in ["bold", "semibold"]:
                font.bold = True
            if "style" in font_config and font_config["style"] == "italic":
                font.italic = True

        # Color
        color_hex = styling.get("color")
        if color_hex and color_hex.startswith("#") and paragraph.runs:
            rgb = self._hex_to_rgb(color_hex)
            paragraph.runs[0].font.color.rgb = RGBColor(*rgb)

        # Alignment
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }
        alignment = styling.get("alignment", "left")
        paragraph.alignment = alignment_map.get(alignment, PP_ALIGN.LEFT)

        # Spacing
        spacing = styling.get("spacing", {})
        if spacing:
            line_height = spacing.get("line_height", 1.2)
            paragraph.line_spacing = line_height

    async def _render_image_element(
        self,
        slide,
        element: dict[str, Any],
        image_uri: str,
        position: dict[str, Any],
        styling: dict[str, Any],
    ) -> None:
        """
        Render an image element.

        Args:
            slide: Slide object
            element: Full element configuration
            image_uri: URI of the image (from content_map)
            position: Position configuration
            styling: Styling configuration
        """
        try:
            # Get position
            left = Inches(position.get("x", 1.0))
            top = Inches(position.get("y", 2.0))
            width = Inches(position.get("width", 4.0))
            height = Inches(position.get("height", 3.0))

            # Download image if it's a URL
            if image_uri.startswith(('http://', 'https://')):
                if self._http_client is None:
                    self._http_client = httpx.AsyncClient(timeout=30.0)

                response = await self._http_client.get(image_uri)
                response.raise_for_status()
                image_stream = io.BytesIO(response.content)
                slide.shapes.add_picture(image_stream, left, top, width, height)
            else:
                # Assume it's a local path
                slide.shapes.add_picture(image_uri, left, top, width, height)

            logger.info(f"Successfully added image: {image_uri[:50]}")

        except Exception as e:
            logger.error(f"Failed to add image {image_uri}: {e}")
            # Add placeholder text box instead
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"[Image: {image_uri[:30]}...]"

    async def _render_table_element(
        self,
        slide,
        element: dict[str, Any],
        table_data: str,
        position: dict[str, Any],
        styling: dict[str, Any],
    ) -> None:
        """
        Render a table element.

        Args:
            slide: Slide object
            element: Full element configuration
            table_data: Table data (CSV, TSV, or JSON string)
            position: Position configuration
            styling: Styling configuration
        """
        try:
            # Parse table data
            rows = self._parse_table_data(table_data)

            if not rows or len(rows) < 1:
                logger.warning("No table data to render")
                return

            # Get position
            left = Inches(position.get("x", 1.0))
            top = Inches(position.get("y", 2.0))
            width = Inches(position.get("width", 8.0))
            height = Inches(position.get("height", 3.0))

            # Create table
            num_rows = len(rows)
            num_cols = len(rows[0]) if rows else 0

            table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
            table = table_shape.table

            # Populate table
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)

                    # Apply header styling to first row
                    if row_idx == 0:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
                                run.font.size = Pt(12)

            logger.info(f"Successfully added table: {num_rows}x{num_cols}")

        except Exception as e:
            logger.error(f"Failed to add table: {e}")

    def _parse_table_data(self, table_data: str) -> list[list[str]]:
        """
        Parse table data from various formats.

        Args:
            table_data: Table data string (CSV, TSV, or simple format)

        Returns:
            List of rows, each row is a list of cell values
        """
        rows = []

        # Try to detect format and parse
        lines = table_data.strip().split('\n')

        for line in lines:
            if not line.strip():
                continue

            # Try different delimiters
            if '\t' in line:
                cells = line.split('\t')
            elif '|' in line:
                cells = [c.strip() for c in line.split('|') if c.strip()]
            elif ',' in line:
                cells = [c.strip() for c in line.split(',')]
            else:
                cells = [line]

            rows.append(cells)

        return rows

    def _render_shape_element(
        self,
        slide,
        element: dict[str, Any],
        position: dict[str, Any],
        styling: dict[str, Any],
    ) -> None:
        """
        Render a shape element (rectangle, circle, etc.).

        Args:
            slide: Slide object
            element: Full element configuration
            position: Position configuration
            styling: Styling configuration
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE

            # Get position
            left = Inches(position.get("x", 1.0))
            top = Inches(position.get("y", 2.0))
            width = Inches(position.get("width", 2.0))
            height = Inches(position.get("height", 1.0))

            # Add rectangle by default
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )

            # Apply background color
            background = styling.get("background", "#4472C4")
            fill = shape.fill
            fill.solid()
            rgb = self._hex_to_rgb(background)
            fill.fore_color.rgb = RGBColor(*rgb)

            logger.info("Successfully added shape")

        except Exception as e:
            logger.error(f"Failed to add shape: {e}")
